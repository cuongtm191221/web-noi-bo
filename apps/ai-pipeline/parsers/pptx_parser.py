import io
from .base import ParsedDocument, PageContent


def parse_pptx(file_bytes: bytes) -> ParsedDocument:
    """Parse PPTX with python-pptx."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    pages = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        pages.append(PageContent(
            text="\n".join(texts),
            location={"slide_number": i + 1},
        ))
    return ParsedDocument(
        pages=pages,
        metadata={"format": "pptx", "slide_count": len(prs.slides)},
    )